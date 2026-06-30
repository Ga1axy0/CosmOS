#!/usr/bin/env python3
"""
Generate a concise project showcase / defense deck for CosmOS.

The script intentionally keeps the deck editable: most diagrams are native
PowerPoint shapes, while selected documentation figures are converted to PNG
and embedded for readability.
"""

from __future__ import annotations

import os
import subprocess
import sys
from pathlib import Path

LOCAL_DEPS = Path("/tmp/xxos_ppt_deps")
if LOCAL_DEPS.exists():
    sys.path.insert(0, str(LOCAL_DEPS))

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ModuleNotFoundError as exc:
    raise SystemExit(
        "Missing python-pptx. Install it with:\n"
        "  python3 -m pip install --target /tmp/xxos_ppt_deps python-pptx"
    ) from exc


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "cosmos_docs" / "assets"
OUT_DIR = ROOT / "docs" / "ppt_assets"
OUT_FILE = ROOT / "cosmos_docs" / "CosmOS_项目特色展示与答辩.pptx"

INKSCAPE = Path("/mnt/c/Program Files/Inkscape/bin/inkscape.exe")

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BLUE = RGBColor(15, 82, 186)
BLUE_DARK = RGBColor(8, 39, 94)
BLUE_LIGHT = RGBColor(228, 238, 255)
RED = RGBColor(210, 48, 48)
RED_LIGHT = RGBColor(255, 235, 235)
TEXT = RGBColor(35, 45, 60)
MUTED = RGBColor(100, 112, 130)
BG = RGBColor(247, 250, 255)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(210, 220, 235)
GREEN = RGBColor(43, 145, 120)
YELLOW = RGBColor(246, 183, 64)

FONT = "Microsoft YaHei"
FONT_LATIN = "Aptos"


def windows_path(path: Path) -> str:
    return subprocess.check_output(["wslpath", "-w", str(path)], text=True).strip()


def convert_svg(name: str, width: int = 1800) -> Path:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    src = ASSET_DIR / name
    dst = OUT_DIR / (src.stem.replace(".", "_") + ".png")
    if dst.exists() and dst.stat().st_mtime >= src.stat().st_mtime:
        return dst
    if not INKSCAPE.exists():
        raise RuntimeError("Inkscape is not available")
    subprocess.run(
        [
            str(INKSCAPE),
            windows_path(src),
            "--export-type=png",
            f"--export-filename={windows_path(dst)}",
            f"--export-width={width}",
        ],
        check=True,
    )
    return dst


def rgb(value: RGBColor) -> str:
    return f"{value[0]:02X}{value[1]:02X}{value[2]:02X}"


def set_fill(shape, color: RGBColor, transparency: int = 0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color: RGBColor = LINE, width: float = 1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def no_line(shape):
    shape.line.fill.background()


def add_text(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    size: int = 18,
    color: RGBColor = TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = FONT,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str | None = None, section: str | None = None):
    if section:
        add_text(slide, section, Inches(0.62), Inches(0.25), Inches(2.6), Inches(0.25), 9, RED, True)
    add_text(slide, title, Inches(0.62), Inches(0.48), Inches(8.8), Inches(0.48), 25, BLUE_DARK, True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.64), Inches(0.98), Inches(9.8), Inches(0.34), 11, MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.29), Inches(1.25), Inches(0.04))
    set_fill(line, RED)
    no_line(line)


def add_footer(slide, index: int):
    add_text(slide, "CosmOS · 项目特色展示与汇报答辩", Inches(0.62), Inches(7.12), Inches(4.8), Inches(0.22), 8, MUTED)
    add_text(slide, f"{index:02d}", Inches(12.15), Inches(7.08), Inches(0.55), Inches(0.24), 9, MUTED, True, PP_ALIGN.RIGHT)


def blank_slide(prs, index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    add_footer(slide, index)
    return slide


def card(slide, x, y, w, h, title: str | None = None, fill: RGBColor = WHITE, line: RGBColor = LINE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(shape, fill)
    set_line(shape, line, 1)
    shape.adjustments[0] = 0.08
    if title:
        add_text(slide, title, x + Inches(0.18), y + Inches(0.14), w - Inches(0.36), Inches(0.28), 13, BLUE_DARK, True)
    return shape


def pill(slide, text: str, x, y, w, color: RGBColor = BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.32))
    set_fill(shape, color)
    no_line(shape)
    shape.adjustments[0] = 0.5
    add_text(slide, text, x, y + Inches(0.065), w, Inches(0.15), 8, WHITE, True, PP_ALIGN.CENTER)
    return shape


def bullets(slide, items: list[str], x, y, w, h, size: int = 16, color: RGBColor = TEXT, gap: float = 0.08):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8 + gap * 10)
        run = p.add_run()
        run.text = item
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def image_contain(slide, img: Path, x, y, w, h):
    from PIL import Image

    with Image.open(img) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    return slide.shapes.add_picture(str(img), px, py, width=pw, height=ph)


def connector(slide, x1, y1, x2, y2, color: RGBColor = BLUE, width: float = 1.6):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def metric_card(slide, x, y, w, h, number: str, label: str, accent: RGBColor = RED):
    card(slide, x, y, w, h)
    add_text(slide, number, x + Inches(0.16), y + Inches(0.18), w - Inches(0.32), Inches(0.42), 24, accent, True)
    add_text(slide, label, x + Inches(0.18), y + Inches(0.68), w - Inches(0.36), Inches(0.45), 10, MUTED)


def bar_chart(slide, x, y, w, h, title: str, rows: list[tuple[str, float, float, str]]):
    card(slide, x, y, w, h, title)
    maxv = max(max(a, b) for _, a, b, _ in rows)
    base_x = x + Inches(1.55)
    bar_w = w - Inches(2.25)
    row_h = (h - Inches(0.82)) / len(rows)
    for i, (label, before, after, note) in enumerate(rows):
        yy = y + Inches(0.72) + row_h * i
        add_text(slide, label, x + Inches(0.18), yy + Inches(0.03), Inches(1.18), Inches(0.18), 8, TEXT, True)
        b1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, base_x, yy, bar_w * before / maxv, Inches(0.11))
        set_fill(b1, RGBColor(170, 185, 205))
        no_line(b1)
        b2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, base_x, yy + Inches(0.17), bar_w * after / maxv, Inches(0.11))
        set_fill(b2, BLUE)
        no_line(b2)
        add_text(slide, note, base_x + bar_w + Inches(0.08), yy + Inches(0.08), Inches(0.45), Inches(0.16), 8, RED, True)


def tag_row(slide, tags: list[str], x, y, colors: list[RGBColor] | None = None):
    cx = x
    colors = colors or [BLUE] * len(tags)
    for tag, color in zip(tags, colors):
        width = Inches(0.2 + max(0.65, len(tag) * 0.105))
        pill(slide, tag, cx, y, width, color)
        cx += width + Inches(0.12)


def make_deck():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "CosmOS 项目特色展示与汇报答辩"
    prs.core_properties.subject = "Operating system project showcase"
    prs.core_properties.author = "CosmOS Team"

    assets = {
        "overview": convert_svg("ch1-summary.svg"),
        "mm": convert_svg("mm_overview.svg"),
        "fs_stack": convert_svg("filesystem_stack_bluegreen.drawio.svg"),
        "poll": convert_svg("poll_bitmap.svg"),
        "net": convert_svg("net_stack.svg"),
        "net_poll": convert_svg("net_poll_wakeup.svg"),
        "hal": convert_svg("hal_pagingarch.svg"),
        "agile": convert_svg("perf_probe_agile_loop.svg"),
        "treemap": convert_svg("kernel_code_treemap.svg"),
        "page_cache": convert_svg("fs_page_cache_iozone.svg"),
        "fastpath": convert_svg("fs_fastpath_iozone.svg"),
        "scheduler": ASSET_DIR / "scheduler_path_bluegreen.png",
        "contributors": ASSET_DIR / "Contributors.png",
    }

    i = 1
    slide = blank_slide(prs, i)
    add_text(slide, "CosmOS", Inches(0.75), Inches(0.72), Inches(4.1), Inches(0.75), 46, BLUE_DARK, True, font=FONT_LATIN)
    add_text(slide, "项目特色展示与汇报答辩", Inches(0.82), Inches(1.56), Inches(5.7), Inches(0.45), 25, TEXT, True)
    add_text(slide, "从 rCore 骨架到面向真实用户程序的实验型 Linux 兼容内核", Inches(0.86), Inches(2.13), Inches(7.9), Inches(0.3), 14, MUTED)
    tag_row(slide, ["SMP", "Linux 语义", "统一事件机制", "跨架构 HAL", "性能观测"], Inches(0.86), Inches(2.72), [BLUE, BLUE, RED, BLUE, RED])
    card(slide, Inches(7.6), Inches(0.92), Inches(4.7), Inches(4.78), fill=WHITE)
    layers = [
        ("用户程序与 Linux 兼容环境", "进程 / 文件描述符 / 信号 / 套接字", BLUE),
        ("系统服务层", "VFS · Signal · Poll · Network · Virtual FS", RED),
        ("内核核心层", "调度 · 进程 · 内存 · 同步 · 回收", BLUE),
        ("HAL 与平台层", "RISC-V 64 / LoongArch 64", BLUE_DARK),
    ]
    for n, (t, s, c) in enumerate(layers):
        yy = Inches(1.3 + n * 0.95)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.05), yy, Inches(3.85), Inches(0.62))
        set_fill(shape, c if n in (1, 3) else BLUE_LIGHT)
        set_line(shape, c, 1)
        shape.adjustments[0] = 0.08
        add_text(slide, t, Inches(8.22), yy + Inches(0.10), Inches(3.5), Inches(0.18), 11, WHITE if n in (1, 3) else BLUE_DARK, True, PP_ALIGN.CENTER)
        add_text(slide, s, Inches(8.22), yy + Inches(0.34), Inches(3.5), Inches(0.16), 8, WHITE if n in (1, 3) else MUTED, False, PP_ALIGN.CENTER)
    add_text(slide, "答辩主线：用稳定基础设施承载真实 Linux 风格语义", Inches(0.86), Inches(5.72), Inches(10.2), Inches(0.34), 18, RED, True)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "一句话定位", "不是最小教学内核，而是可运行复杂用户程序的实验型 OS", "01 / 项目定位")
    metric_card(slide, Inches(0.78), Inches(1.72), Inches(2.65), Inches(1.28), "55k+", "自主内核非注释代码量", RED)
    metric_card(slide, Inches(3.68), Inches(1.72), Inches(2.65), Inches(1.28), "2", "RISC-V 64 与 LoongArch 64", BLUE)
    metric_card(slide, Inches(6.58), Inches(1.72), Inches(2.65), Inches(1.28), "5级", "文件系统缓存与可旁路实验", RED)
    metric_card(slide, Inches(9.48), Inches(1.72), Inches(2.65), Inches(1.28), "统一", "等待队列 / 信号 / poll / socket", BLUE)
    card(slide, Inches(0.78), Inches(3.42), Inches(11.35), Inches(2.28), "核心目标")
    bullets(
        slide,
        [
            "尽量接近 Linux 的进程、文件、信号、网络和设备语义",
            "在 SMP 环境下维护调度、阻塞唤醒、页表修改和资源回收的不变式",
            "通过 HAL 隔离架构差异，让内核主体复用到两套 ISA",
        ],
        Inches(1.05),
        Inches(4.02),
        Inches(10.6),
        Inches(1.2),
        17,
    )

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "汇报路线", "用 6 个关键词快速说明 CosmOS 的工程价值", "02 / 叙事结构")
    nodes = [
        ("资源模型", "PCB/TCB 分离\nfork/exec/exit/wait 生命周期"),
        ("并发正确性", "per-hart 调度\n三段式阻塞与唤醒"),
        ("内存一致性", "VMA + PageTable\nCOW / TLB shootdown"),
        ("I/O 数据面", "统一 File 对象\n五级缓存与快路径"),
        ("事件与网络", "Signal / Poll / Socket\n软中断式 polling"),
        ("跨架构与观测", "HAL trait\nio_perf / perf_probe"),
    ]
    for idx, (t, s) in enumerate(nodes):
        col, row = idx % 3, idx // 3
        x, y = Inches(0.75 + col * 4.05), Inches(1.75 + row * 2.1)
        card(slide, x, y, Inches(3.45), Inches(1.28), fill=WHITE)
        pill(slide, f"{idx + 1}", x + Inches(0.18), y + Inches(0.18), Inches(0.38), RED if idx in (1, 3) else BLUE)
        add_text(slide, t, x + Inches(0.72), y + Inches(0.17), Inches(2.4), Inches(0.25), 15, BLUE_DARK, True)
        add_text(slide, s, x + Inches(0.28), y + Inches(0.58), Inches(2.85), Inches(0.5), 10, MUTED)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "系统总览", "四层结构 + 横向统一对象模型", "03 / 总体架构")
    image_contain(slide, assets["overview"], Inches(0.68), Inches(1.48), Inches(7.25), Inches(4.95))
    card(slide, Inches(8.25), Inches(1.66), Inches(3.9), Inches(4.42), "需要记住的三条主线")
    bullets(
        slide,
        [
            "统一 File 对象承载普通文件、设备、管道、终端和套接字",
            "等待队列 / poll / 信号贯穿文件、进程、同步和网络",
            "HAL 将 trap、IRQ、paging、timer 差异压到边界",
        ],
        Inches(8.55),
        Inches(2.3),
        Inches(3.28),
        Inches(2.6),
        14,
    )
    add_text(slide, "重点不是模块堆叠，而是共享机制沉淀。", Inches(8.55), Inches(5.35), Inches(3.25), Inches(0.3), 13, RED, True)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "代码组织与工程体量", "模块边界和文档结构基本对齐", "04 / 工程规模")
    image_contain(slide, assets["treemap"], Inches(0.7), Inches(1.45), Inches(6.5), Inches(4.9))
    card(slide, Inches(7.55), Inches(1.55), Inches(4.55), Inches(4.65), "目录映射")
    rows = [
        ("os/src/sched", "调度类、运行队列、上下文切换"),
        ("os/src/task", "PCB / TCB / 进程生命周期"),
        ("os/src/mm", "页表、VMA、COW、TLB shootdown"),
        ("os/src/fs + fs/", "VFS、多后端、缓存、虚拟文件系统"),
        ("os/src/net", "socket、smoltcp、VirtIO-net"),
        ("os/src/hal + arch + platform", "跨架构抽象与平台适配"),
    ]
    yy = Inches(2.1)
    for path, desc in rows:
        add_text(slide, path, Inches(7.85), yy, Inches(1.75), Inches(0.18), 8, RED, True, font=FONT_LATIN)
        add_text(slide, desc, Inches(9.55), yy, Inches(2.25), Inches(0.18), 9, TEXT)
        yy += Inches(0.48)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "特色一：SMP 调度不是一个队列", "per-hart 运行队列 + 多调度类 + 延迟抢占", "05 / 调度")
    image_contain(slide, assets["scheduler"], Inches(0.72), Inches(1.48), Inches(7.0), Inches(4.9))
    card(slide, Inches(8.05), Inches(1.58), Inches(4.05), Inches(4.65), "设计重点")
    bullets(
        slide,
        [
            "实时类：SCHED_FIFO / SCHED_RR，优先级语义明确",
            "普通类：CFS 风格 vruntime，nice 权重影响公平时间轴",
            "timer / IPI 只设置 resched_reason，在安全点统一切换",
            "任务先切回 hart 的 idle 调度上下文，再选择下一任务",
        ],
        Inches(8.35),
        Inches(2.15),
        Inches(3.35),
        Inches(2.9),
        13,
    )
    pill(slide, "核心价值：减少锁竞争，并让上下文切换边界可推理", Inches(8.35), Inches(5.62), Inches(3.45), RED)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "阻塞—唤醒正确性", "task_status / on_rq / on_cpu 三个状态各管一件事", "06 / SMP 不变式")
    names = [
        ("task_status", "语义状态\nRunnable / Interruptible / Zombie"),
        ("sched.on_rq", "是否已挂入\n某条运行队列"),
        ("on_cpu", "寄存器现场是否\n仍归某个 hart 所有"),
    ]
    for n, (t, s) in enumerate(names):
        x = Inches(0.9 + n * 4.1)
        card(slide, x, Inches(1.7), Inches(3.15), Inches(1.55), t, fill=BLUE_LIGHT if n != 2 else RED_LIGHT, line=BLUE if n != 2 else RED)
        add_text(slide, s, x + Inches(0.28), Inches(2.28), Inches(2.55), Inches(0.48), 13, TEXT, False, PP_ALIGN.CENTER)
    connector(slide, Inches(2.48), Inches(3.35), Inches(2.48), Inches(4.05), RED, 1.8)
    connector(slide, Inches(6.58), Inches(3.35), Inches(6.58), Inches(4.05), RED, 1.8)
    connector(slide, Inches(10.68), Inches(3.35), Inches(10.68), Inches(4.05), RED, 1.8)
    card(slide, Inches(1.15), Inches(4.1), Inches(10.9), Inches(1.25), "关键协议")
    bullets(
        slide,
        [
            "阻塞路径：入等待队列 → 复查条件 → 真正切出 CPU，关闭丢失唤醒窗口",
            "跨 hart 唤醒：若 on_cpu 仍为 true，等待原 hart 完成寄存器保存后再入队",
            "退出路径：当前任务内核栈延迟释放，避免切换仍踩在已回收栈上",
        ],
        Inches(1.45),
        Inches(4.62),
        Inches(10.2),
        Inches(0.55),
        12,
    )
    add_text(slide, "结论：任务不会同时处在 current 与 runqueue，也不会在半保存现场被迁移。", Inches(1.16), Inches(5.75), Inches(10.8), Inches(0.25), 14, RED, True, PP_ALIGN.CENTER)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "特色二：进程是资源容器，任务是调度实体", "fork / clone / exec / exit / wait 围绕 PCB 生命周期展开", "07 / 进程管理")
    stages = [
        ("clone/fork", "复制或共享资源\n修补 trap context"),
        ("exec", "保留 PID\n替换用户态映像"),
        ("exit", "进入 zombie\n释放可立即回收资源"),
        ("wait", "父进程读取状态\n真正 drop PCB"),
    ]
    for idx, (t, s) in enumerate(stages):
        x = Inches(0.85 + idx * 3.05)
        card(slide, x, Inches(1.82), Inches(2.55), Inches(1.2), t, fill=WHITE)
        add_text(slide, s, x + Inches(0.2), Inches(2.34), Inches(2.15), Inches(0.38), 11, MUTED, False, PP_ALIGN.CENTER)
        if idx < len(stages) - 1:
            connector(slide, x + Inches(2.58), Inches(2.42), x + Inches(2.98), Inches(2.42), BLUE)
    card(slide, Inches(1.0), Inches(3.75), Inches(5.25), Inches(1.55), "PCB 管资源")
    bullets(slide, ["地址空间 / fd 表 / cwd-root / 凭据", "信号处置 / 资源限制 / 子进程集合", "线程集合 / 时间统计 / wait_exit_queue"], Inches(1.28), Inches(4.25), Inches(4.7), Inches(0.85), 12)
    card(slide, Inches(7.0), Inches(3.75), Inches(5.25), Inches(1.55), "TCB 管执行")
    bullets(slide, ["内核栈 / trap context / 调度状态", "Linux 可见 tid / clear_child_tid", "等待原因 / 信号屏蔽 / on_cpu"], Inches(7.28), Inches(4.25), Inches(4.7), Inches(0.85), 12)
    add_text(slide, "发布顺序是红线：用户栈、trap context、PID 表项准备好之后，任务才能入队。", Inches(1.05), Inches(5.92), Inches(10.9), Inches(0.25), 13, RED, True, PP_ALIGN.CENTER)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "特色三：内存管理是全内核的数据面", "VMA 保存语义，PageTable 承载硬件状态", "08 / 内存管理")
    image_contain(slide, assets["mm"], Inches(0.68), Inches(1.45), Inches(7.05), Inches(4.95))
    card(slide, Inches(8.05), Inches(1.62), Inches(4.05), Inches(4.62), "三层理解")
    bullets(
        slide,
        [
            "基础资源层：bootinfo → BuddyFrameAllocator → FrameTracker",
            "地址空间层：KERNEL_SPACE 与每进程 MemorySet / VMA",
            "运行期机制：page fault、COW、page cache、TLB shootdown",
        ],
        Inches(8.35),
        Inches(2.25),
        Inches(3.35),
        Inches(2.0),
        13,
    )
    add_text(slide, "旧页不能立即释放：先清 PTE，再 shootdown，最后 deferred reclaim。", Inches(8.35), Inches(5.28), Inches(3.35), Inches(0.44), 12, RED, True)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "VMA / COW / Page Cache / TLB 的一致性", "把用户可见语义落到页生命周期", "09 / 内存关键路径")
    left = [
        ("合法但未映射", "VMA 存在，PTE 不存在\n缺页时按语义物化"),
        ("私有页 / COW", "fork 后降权共享\n写缺页再复制"),
        ("文件映射", "MAP_SHARED 可直接映射 page cache\ntruncate 需要失效映射"),
        ("回收与 shootdown", "降权/拆映射/换页\n必须处理远端 TLB"),
    ]
    for idx, (t, s) in enumerate(left):
        x = Inches(0.8 + (idx % 2) * 6.05)
        y = Inches(1.65 + (idx // 2) * 2.05)
        card(slide, x, y, Inches(5.25), Inches(1.18), t, fill=WHITE)
        add_text(slide, s, x + Inches(0.3), y + Inches(0.54), Inches(4.65), Inches(0.45), 11, MUTED)
    add_text(slide, "一句话：页表只是硬件事实，VMA 才是可恢复的用户语义。", Inches(1.2), Inches(6.0), Inches(10.7), Inches(0.25), 16, RED, True, PP_ALIGN.CENTER)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "特色四：统一文件对象与五级缓存", "普通文件、设备、pipe、tty、socket 共享 fd / poll 入口", "10 / 文件子系统")
    image_contain(slide, assets["fs_stack"], Inches(0.55), Inches(1.36), Inches(7.4), Inches(5.15))
    card(slide, Inches(8.25), Inches(1.55), Inches(3.95), Inches(4.72), "缓存分工")
    bullets(
        slide,
        [
            "Stat：inode 属性快照",
            "Dentry：路径分量解析",
            "Inode：同一文件身份复用",
            "Page：4 KB 文件数据页",
            "Block：512 B 磁盘块收口",
        ],
        Inches(8.55),
        Inches(2.18),
        Inches(3.2),
        Inches(2.6),
        13,
    )
    add_text(slide, "独立旁路开关让性能收益可以逐级归因。", Inches(8.55), Inches(5.55), Inches(3.2), Inches(0.28), 12, RED, True)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "文件 I/O 性能：缓存收益清晰可解释", "iozone / 目录遍历 / du / ls-al 均可用计数器交叉验证", "11 / 性能数据")
    bar_chart(
        slide,
        Inches(0.78),
        Inches(1.55),
        Inches(5.5),
        Inches(4.75),
        "Page cache 对 iozone 的影响（KB/s）",
        [
            ("顺序读", 4943, 111371, "22x"),
            ("顺序写", 2558, 88171, "34x"),
            ("随机读", 1.0, 3.0, "3x"),
            ("随机写", 1.0, 2.2, "2.2x"),
        ],
    )
    image_contain(slide, assets["page_cache"], Inches(6.55), Inches(1.55), Inches(5.75), Inches(4.75))
    add_text(slide, "读写提升不是“凭感觉”：/proc/io_perf 能看到命中、未命中、回写与淘汰。", Inches(1.05), Inches(6.45), Inches(11.0), Inches(0.24), 13, RED, True, PP_ALIGN.CENTER)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "小 I/O 快路径：优化每次 syscall 的固定成本", "fd 查找缓存 + 单页零分配缓冲 + 每描述符热页缓存", "12 / I/O 快路径")
    image_contain(slide, assets["fastpath"], Inches(0.68), Inches(1.35), Inches(7.45), Inches(5.0))
    card(slide, Inches(8.45), Inches(1.68), Inches(3.7), Inches(4.25), "收益特征")
    bullets(
        slide,
        [
            "1 KB / 4 KB 小记录普遍提升 20%–30%",
            "4 KB 随机读提升约 3.3 倍",
            "64 KB / 256 KB 大记录收益回落",
            "这是 syscall 固定成本优化，不改变语义路径",
        ],
        Inches(8.75),
        Inches(2.25),
        Inches(3.05),
        Inches(2.25),
        13,
    )
    add_text(slide, "所有快路径都有旁路开关，便于回归与归因。", Inches(8.75), Inches(5.42), Inches(3.05), Inches(0.3), 12, RED, True)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "特色五：统一阻塞—唤醒基础设施", "slot registry + generation + keyed wait queue", "13 / 事件机制")
    image_contain(slide, assets["poll"], Inches(0.75), Inches(1.48), Inches(6.15), Inches(4.8))
    card(slide, Inches(7.25), Inches(1.52), Inches(4.95), Inches(4.75), "同一种范式承载四类等待")
    rows = [
        ("ppoll / pselect6", "二维位图注册表，fd 行 × poll key 列"),
        ("FUTEX_WAIT", "按用户地址注册，支持 requeue"),
        ("rt_sigtimedwait", "同步等待并消费信号"),
        ("socket timeout", "SO_RCVTIMEO / SO_SNDTIMEO 超时唤醒"),
    ]
    yy = Inches(2.12)
    for t, s in rows:
        pill(slide, t, Inches(7.55), yy, Inches(1.55), BLUE)
        add_text(slide, s, Inches(9.3), yy + Inches(0.05), Inches(2.4), Inches(0.18), 9, TEXT)
        yy += Inches(0.72)
    add_text(slide, "代计数器防 ABA，锁内收集、锁外唤醒，关闭丢失唤醒窗口。", Inches(7.55), Inches(5.55), Inches(4.05), Inches(0.3), 12, RED, True)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "信号：异步事件如何接回同步阻塞", "终端 Ctrl+C、EINTR、SA_RESTART 都落在同一条路径上", "14 / Signal")
    steps = [
        ("事件到达", "tty / kill / timer"),
        ("置 pending", "进程级或线程级信号位"),
        ("摘队唤醒", "current_wq_handle / poll key"),
        ("返回用户态", "构造 signal frame"),
        ("sigreturn", "恢复上下文或重启 syscall"),
    ]
    for idx, (t, s) in enumerate(steps):
        x = Inches(0.65 + idx * 2.48)
        card(slide, x, Inches(2.05), Inches(2.05), Inches(1.05), t, fill=RED_LIGHT if idx in (2, 4) else WHITE, line=RED if idx in (2, 4) else LINE)
        add_text(slide, s, x + Inches(0.16), Inches(2.55), Inches(1.72), Inches(0.28), 9, MUTED, False, PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            connector(slide, x + Inches(2.05), Inches(2.58), x + Inches(2.42), Inches(2.58), BLUE)
    card(slide, Inches(1.0), Inches(4.15), Inches(5.1), Inches(1.2), "EINTR 判定")
    add_text(slide, "只有“确实会被采取行动”的信号才打断阻塞调用；默认忽略或被屏蔽的信号只挂起。", Inches(1.28), Inches(4.7), Inches(4.55), Inches(0.36), 12, TEXT)
    card(slide, Inches(7.0), Inches(4.15), Inches(5.1), Inches(1.2), "SA_RESTART")
    add_text(slide, "保存 orig_a0 与 syscall PC；处理函数返回后回退到 syscall 指令重新执行。", Inches(7.28), Inches(4.7), Inches(4.55), Inches(0.36), 12, TEXT)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "特色六：网络栈选择成熟协议核心，自己做好驱动模型", "smoltcp + VirtIO-net + socket 文件对象", "15 / 网络栈")
    image_contain(slide, assets["net"], Inches(0.68), Inches(1.43), Inches(7.35), Inches(5.05))
    card(slide, Inches(8.35), Inches(1.62), Inches(3.8), Inches(4.55), "关键设计")
    bullets(
        slide,
        [
            "smoltcp 协作式推进：内核决定何时 poll",
            "socket 实现 File trait，天然接入 fd 与 poll",
            "read_wait / write_wait 复用等待队列",
            "AF_UNIX / Raw IPv6 / netlink 兼容层服务真实工具",
        ],
        Inches(8.65),
        Inches(2.22),
        Inches(3.05),
        Inches(2.65),
        13,
    )
    add_text(slide, "协议实现成熟，工程重点转向事件驱动与阻塞语义。", Inches(8.65), Inches(5.48), Inches(3.05), Inches(0.3), 12, RED, True)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "软中断式 Polling：中断只设标志，重活统一推进", "NEED_POLL + poll_at 构成双触发模型", "16 / 网络事件")
    image_contain(slide, assets["net_poll"], Inches(0.62), Inches(1.48), Inches(7.55), Inches(4.9))
    card(slide, Inches(8.45), Inches(1.74), Inches(3.72), Inches(3.95), "为什么这样做")
    bullets(
        slide,
        [
            "硬中断路径短：应答设备、置位 NEED_POLL",
            "定时器按 smoltcp poll_at 推进 TCP 时间事件",
            "poll_once 后统一唤醒 socket wait 与 poll registry",
            "预算机制平衡低延迟和不霸占 CPU",
        ],
        Inches(8.75),
        Inches(2.3),
        Inches(3.02),
        Inches(2.35),
        12,
    )

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "跨架构 HAL：把硬件事实压到 trait 边界", "内核主体不直接感知 RISC-V 还是 LoongArch", "17 / HAL")
    image_contain(slide, assets["hal"], Inches(0.65), Inches(1.42), Inches(7.75), Inches(5.05))
    card(slide, Inches(8.72), Inches(1.6), Inches(3.45), Inches(4.6), "封装差异")
    bullets(
        slide,
        [
            "TrapMachine：陷阱原因、返回用户态、syscall 长度",
            "PagingArch：PTE 编码、token、TLB、页表索引",
            "SignalAbi：ucontext / rt_sigaction 字节布局",
            "Platform：timer、IPI、SMP 启动、中断控制器",
        ],
        Inches(9.0),
        Inches(2.2),
        Inches(2.82),
        Inches(2.5),
        11,
    )
    add_text(slide, "LoongArch 支持验证了抽象边界：主体代码复用，移植工作有界。", Inches(8.98), Inches(5.42), Inches(2.95), Inches(0.42), 12, RED, True)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "Linux 兼容面：优先覆盖用户程序依赖的关键语义", "不是追求全量复制，而是把常用路径做深做对", "18 / 兼容性")
    cols = [
        ("进程", ["fork / clone / execve", "wait4 / zombie / reparent", "进程组 / 会话 / 凭据"]),
        ("文件", ["open / dup / fcntl", "procfs / tmpfs / devfs", "pipe / tty / symlink"]),
        ("内存", ["mmap / munmap / brk", "COW / lazy allocation", "file-backed mmap"]),
        ("事件", ["signal / sigreturn", "futex / ppoll / pselect6", "SA_RESTART / EINTR"]),
        ("网络", ["TCP / UDP sockets", "AF_UNIX / SCM_RIGHTS", "netlink / packet / AF_ALG"]),
    ]
    for idx, (t, items) in enumerate(cols):
        x = Inches(0.55 + idx * 2.5)
        card(slide, x, Inches(1.65), Inches(2.15), Inches(3.95), t, fill=WHITE)
        yy = Inches(2.32)
        for item in items:
            add_text(slide, "■", x + Inches(0.18), yy, Inches(0.15), Inches(0.15), 8, RED, True)
            add_text(slide, item, x + Inches(0.42), yy - Inches(0.01), Inches(1.45), Inches(0.24), 9, TEXT)
            yy += Inches(0.62)
    add_text(slide, "答辩要点：兼容性来自对象模型与生命周期，而不是 syscall 清单堆叠。", Inches(1.0), Inches(6.18), Inches(11.0), Inches(0.25), 14, RED, True, PP_ALIGN.CENTER)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "敏捷性能观测：优化必须可解释", "/proc/io_perf 看事件，perf_probe 看耗时", "19 / 工具链")
    image_contain(slide, assets["agile"], Inches(0.72), Inches(1.45), Inches(6.55), Inches(4.85))
    card(slide, Inches(7.6), Inches(1.62), Inches(4.55), Inches(4.45), "两类工具")
    metric_card(slide, Inches(7.95), Inches(2.15), Inches(1.72), Inches(1.15), "io_perf", "长期语义事件计数", BLUE)
    metric_card(slide, Inches(10.05), Inches(2.15), Inches(1.72), Inches(1.15), "probe", "短期命名耗时探针", RED)
    bullets(
        slide,
        [
            "热路径只做原子计数，读取 procfs 时再渲染",
            "probe 可运行期开关与 reset，正式性能用无探针构建确认",
            "每轮围绕一个假设插桩、复现、读数、保留或回退",
        ],
        Inches(7.95),
        Inches(3.68),
        Inches(3.75),
        Inches(1.35),
        12,
    )

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "协作方式：专项分工 + 跨模块共同收敛", "功能补齐、专项重构、测试调优三个阶段", "20 / 团队协作")
    image_contain(slide, assets["contributors"], Inches(0.7), Inches(1.42), Inches(5.3), Inches(4.65))
    stages = [
        ("早期", "Basic Test 驱动\n快速补齐骨架"),
        ("中期", "调度/进程/FS/MM/Net/HAL\n长期方向分工"),
        ("后期", "性能瓶颈、并发竞态\n跨架构差异共同调试"),
    ]
    for idx, (t, s) in enumerate(stages):
        x = Inches(6.55)
        y = Inches(1.7 + idx * 1.38)
        card(slide, x, y, Inches(5.3), Inches(0.92), t, fill=BLUE_LIGHT if idx != 2 else RED_LIGHT, line=BLUE if idx != 2 else RED)
        add_text(slide, s, x + Inches(1.0), y + Inches(0.18), Inches(3.7), Inches(0.4), 12, TEXT)
    add_text(slide, "自动化测试和可复现实验让主分支始终保持基本可用。", Inches(6.65), Inches(5.72), Inches(5.05), Inches(0.26), 13, RED, True)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "诚实的边界：当前限制与取舍", "实验内核走向真实负载时必须面对的成本", "21 / 不足")
    limits = [
        ("兼容性覆盖", "部分 ioctl、socket option、特殊文件语义仍需补齐"),
        ("缓存分层", "page cache 与 block cache 存在文件数据双重缓存"),
        ("网络扩展性", "单一 NET_STACK 锁与逐帧拷贝限制多连接吞吐"),
        ("并发长尾", "极端时序下仍需更多压力测试和不变式检查"),
        ("评估体系", "benchmark 尚需长期化、自动化沉淀"),
    ]
    for idx, (t, s) in enumerate(limits):
        x = Inches(0.82 + (idx % 2) * 5.95)
        y = Inches(1.58 + (idx // 2) * 1.42)
        w = Inches(5.15) if idx < 4 else Inches(11.1)
        card(slide, x if idx < 4 else Inches(0.82), y, w, Inches(0.9), t, fill=RED_LIGHT if idx in (1, 2) else WHITE, line=RED if idx in (1, 2) else LINE)
        add_text(slide, s, (x if idx < 4 else Inches(0.82)) + Inches(0.28), y + Inches(0.42), w - Inches(0.55), Inches(0.2), 10, TEXT)
    add_text(slide, "取舍原则：先保证结构清晰、语义可解释、测试可通过，再逐步降成本。", Inches(1.0), Inches(6.16), Inches(11.0), Inches(0.25), 14, RED, True, PP_ALIGN.CENTER)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "后续路线", "先补语义，再降成本，最后扩展规模", "22 / 展望")
    roadmap = [
        ("近期", "补齐 Linux 兼容缺口\n清理临时分支\n增加回归测试", BLUE),
        ("中期", "统一缓存与 I/O 热路径\n降低全局锁竞争\n沉淀 benchmark 流程", RED),
        ("长期", "扩展平台与设备支持\n系统化验证并发不变式\n形成完整实验平台", BLUE_DARK),
    ]
    for idx, (t, s, c) in enumerate(roadmap):
        x = Inches(0.95 + idx * 4.05)
        card(slide, x, Inches(2.05), Inches(3.3), Inches(2.5), t, fill=WHITE, line=c)
        pill(slide, t, x + Inches(0.35), Inches(2.42), Inches(0.78), c)
        add_text(slide, s, x + Inches(0.42), Inches(3.05), Inches(2.5), Inches(0.85), 13, TEXT)
    add_text(slide, "目标不是继续堆功能，而是在更复杂负载和更强并发下保持一致、稳定、可解释。", Inches(1.12), Inches(5.72), Inches(10.95), Inches(0.3), 14, RED, True, PP_ALIGN.CENTER)

    i += 1
    slide = blank_slide(prs, i)
    add_title(slide, "总结：CosmOS 的核心特色", "用共享机制管理复杂性，用观测闭环约束优化", "23 / 总结")
    summary = [
        ("接近 Linux 的用户程序运行环境", "进程、文件、信号、网络、设备语义协同"),
        ("SMP 下可推理的执行与唤醒协议", "per-hart 调度、三段式阻塞、on_cpu 发布语义"),
        ("贯穿内核的数据面与事件面", "VMA/PageTable/TLB，File/Poll/Signal/Socket"),
        ("跨架构复用和性能可解释", "HAL trait 边界，io_perf 与 perf_probe 闭环"),
    ]
    for idx, (t, s) in enumerate(summary):
        x = Inches(0.85 + (idx % 2) * 5.95)
        y = Inches(1.65 + (idx // 2) * 1.75)
        card(slide, x, y, Inches(5.2), Inches(1.2), t, fill=BLUE_LIGHT if idx % 2 == 0 else WHITE, line=BLUE if idx % 2 == 0 else RED)
        add_text(slide, s, x + Inches(0.28), y + Inches(0.56), Inches(4.6), Inches(0.3), 11, TEXT)
    add_text(slide, "一句话收束：CosmOS 已从教学骨架演进为一个多子系统协同的 OS 实验平台。", Inches(0.95), Inches(6.0), Inches(11.3), Inches(0.3), 15, RED, True, PP_ALIGN.CENTER)

    i += 1
    slide = blank_slide(prs, i)
    add_text(slide, "Q&A", Inches(0.9), Inches(1.2), Inches(4.8), Inches(0.7), 48, BLUE_DARK, True, font=FONT_LATIN)
    add_text(slide, "感谢聆听", Inches(0.96), Inches(2.08), Inches(3.0), Inches(0.42), 25, TEXT, True)
    add_text(slide, "欢迎围绕调度、内存、文件系统、事件机制、网络栈与 HAL 提问", Inches(1.0), Inches(2.82), Inches(8.8), Inches(0.35), 15, MUTED)
    card(slide, Inches(7.35), Inches(1.28), Inches(4.7), Inches(4.15), fill=WHITE)
    add_text(slide, "答辩备忘", Inches(7.75), Inches(1.78), Inches(2.0), Inches(0.28), 15, RED, True)
    bullets(
        slide,
        [
            "讲设计主线，不背系统调用清单",
            "讲不变式，说明为什么并发路径安全",
            "讲数据归因，说明性能优化如何验证",
            "主动承认边界，给出后续路线",
        ],
        Inches(7.78),
        Inches(2.42),
        Inches(3.6),
        Inches(1.75),
        13,
    )

    prs.save(OUT_FILE)
    print(OUT_FILE)


if __name__ == "__main__":
    make_deck()
